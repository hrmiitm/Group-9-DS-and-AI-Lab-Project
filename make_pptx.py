"""
FraudGuard Web Extension — Slide Generator
Produces webextension.pptx in the same directory.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette ─────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x11, 0x1A)   # near-black navy
BG_CARD      = RGBColor(0x13, 0x1C, 0x2E)   # slightly lighter card
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)   # sky-blue headlines
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)   # emerald — SAFE / positive
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)   # amber   — SUSPICIOUS
ACCENT_RED   = RGBColor(0xF8, 0x71, 0x71)   # red     — FAKE / danger
ACCENT_PURP  = RGBColor(0xA7, 0x8B, 0xFA)   # purple  — v2 brand
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
MUTED        = RGBColor(0x94, 0xA3, 0xB8)   # slate-400 body text
DIVIDER      = RGBColor(0x1E, 0x2D, 0x45)   # subtle border

W = Inches(13.33)   # 16:9 width
H = Inches(7.5)     # 16:9 height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_bg(slide):
    """Fill slide background dark."""
    add_rect(slide, 0, 0, W, H, fill=BG_DARK)


def accent_bar(slide, color=ACCENT_BLUE, height=Inches(0.055)):
    """Top accent bar."""
    add_rect(slide, 0, 0, W, height, fill=color)


def slide_title(slide, title, subtitle=None, title_color=ACCENT_BLUE):
    slide_bg(slide)
    accent_bar(slide)
    add_text(slide, title,
             Inches(0.55), Inches(0.18), Inches(12), Inches(0.7),
             size=32, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.55), Inches(0.88), Inches(11), Inches(0.45),
                 size=16, color=MUTED)
    # thin divider line under title area
    add_rect(slide, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.012),
             fill=DIVIDER)


def bullet_box(slide, items, x, y, w, h,
               bullet="▸", size=14, color=WHITE, gap=4):
    """Add a list of bullet items inside a text box."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(gap)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txBox


def card(slide, x, y, w, h, fill=BG_CARD, line=DIVIDER):
    """Rounded-looking card (rectangle with border)."""
    return add_rect(slide, x, y, w, h, fill=fill, line=line, line_w=Pt(1))


def tag(slide, text, x, y, w=Inches(1.6), h=Inches(0.32), color=ACCENT_BLUE, tsize=11):
    """Pill-shaped label."""
    add_rect(slide, x, y, w, h, fill=color)
    add_text(slide, text, x, y + Pt(2), w, h,
             size=tsize, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)

# Gradient top bar (two overlapping rectangles)
add_rect(sl, 0, 0, W, Inches(0.08), fill=ACCENT_BLUE)
add_rect(sl, 0, Inches(0.08), W, Inches(0.04), fill=ACCENT_PURP)

# Big shield emoji area
add_text(sl, "🛡️", Inches(5.5), Inches(1.1), Inches(2), Inches(1.2),
         size=72, align=PP_ALIGN.CENTER)

# Main title
add_text(sl, "FraudGuard", Inches(1), Inches(2.3), Inches(11.3), Inches(1.1),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(sl, "Chrome Extension — From v1 to v2",
         Inches(1), Inches(3.35), Inches(11.3), Inches(0.6),
         size=22, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Description
add_text(sl,
         "AI-Powered LinkedIn Job Fraud Detection  ·  Real-Time  ·  Multi-Signal",
         Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
         size=15, color=MUTED, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(sl, 0, Inches(7.15), W, Inches(0.35), fill=BG_CARD)
add_text(sl, "Group 9  ·  DS & AI Lab Project",
         Inches(0.3), Inches(7.17), Inches(12.7), Inches(0.3),
         size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "Agenda", "What this presentation covers")

agenda = [
    ("01", "The Problem",              "Why LinkedIn job fraud is dangerous"),
    ("02", "Version 1 — Architecture", "How the first extension worked"),
    ("03", "Version 1 — Data Flow",    "Step-by-step: scrape → ML → Gemini → result"),
    ("04", "MCP Tools We Built (v1)",  "Custom tool framework: BaseTool, ToolRegistry, ToolResult"),
    ("05", "Version 2 — Why Rebuild",  "Deploying backend + model for richer analysis"),
    ("06", "Version 2 — Architecture", "New flow: content script + deployed API"),
    ("07", "The 13 Investigation Tools","Every tool, what it checks, and why"),
    ("08", "Verdict Engine",           "Weighted scoring + LLM final summary"),
    ("09", "v1 vs v2 — Side by Side",  "What changed and why it matters"),
]

col_w = Inches(6.0)
for i, (num, title, desc) in enumerate(agenda):
    row = i % 5
    col = i // 5
    x = Inches(0.45) + col * Inches(6.55)
    y = Inches(1.45) + row * Inches(1.0)
    card(sl, x, y, col_w, Inches(0.82))
    add_text(sl, num, x + Inches(0.12), y + Inches(0.08), Inches(0.5), Inches(0.35),
             size=13, bold=True, color=ACCENT_BLUE)
    add_text(sl, title, x + Inches(0.55), y + Inches(0.06), Inches(5.3), Inches(0.35),
             size=14, bold=True, color=WHITE)
    add_text(sl, desc, x + Inches(0.55), y + Inches(0.42), Inches(5.3), Inches(0.32),
             size=11, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "The Problem", "LinkedIn job fraud is widespread, growing, and hard to detect manually")

# Left: stats
card(sl, Inches(0.45), Inches(1.4), Inches(5.9), Inches(5.6))
add_text(sl, "Why This Matters", Inches(0.7), Inches(1.55), Inches(5.4), Inches(0.45),
         size=16, bold=True, color=ACCENT_BLUE)

stats = [
    ("$2B+",   "lost to employment scams yearly (FTC)"),
    ("1 in 5", "LinkedIn users have seen a fake job post"),
    ("48hrs",  "average time before a scam post is removed"),
    ("0",      "native fraud warnings shown to users by LinkedIn"),
]
for i, (num, label) in enumerate(stats):
    yy = Inches(2.15) + i * Inches(1.12)
    add_rect(sl, Inches(0.7), yy, Inches(1.4), Inches(0.75), fill=BG_DARK, line=ACCENT_RED, line_w=Pt(1.5))
    add_text(sl, num, Inches(0.7), yy + Inches(0.08), Inches(1.4), Inches(0.42),
             size=22, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
    add_text(sl, label, Inches(2.25), yy + Inches(0.18), Inches(3.9), Inches(0.4),
             size=12, color=MUTED)

# Right: fraud types
card(sl, Inches(6.8), Inches(1.4), Inches(6.1), Inches(5.6))
add_text(sl, "Common Fraud Patterns", Inches(7.05), Inches(1.55), Inches(5.6), Inches(0.45),
         size=16, bold=True, color=ACCENT_AMBER)

patterns = [
    ("💸  Advance-Fee Scams",    "Ask you to pay for training kits, equipment, or background checks upfront"),
    ("🎣  Phishing Attacks",     "Lure you into submitting personal info via fake 'onboarding' portals"),
    ("🏢  Ghost Companies",      "Completely fabricated businesses with no online presence"),
    ("💼  Identity Theft",       "Request government ID, bank details, or SSN in the 'application'"),
    ("📈  Unrealistic Offers",   "\"$150/hr, no experience needed, work from home, immediate start\""),
]
for i, (title, desc) in enumerate(patterns):
    yy = Inches(2.1) + i * Inches(0.98)
    add_text(sl, title, Inches(7.1), yy, Inches(5.6), Inches(0.32),
             size=13, bold=True, color=WHITE)
    add_text(sl, desc, Inches(7.1), yy + Inches(0.32), Inches(5.7), Inches(0.5),
             size=11, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — V1 ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "Version 1 — Architecture",
            "A self-contained browser extension using a LangChain-inspired tool pipeline")

tag(sl, "V1", Inches(12.3), Inches(0.2), w=Inches(0.75), h=Inches(0.3), color=ACCENT_GREEN)

# Flow boxes
boxes = [
    (Inches(0.4),  "LinkedIn\nJob Page",  MUTED,        "User opens job listing on linkedin.com"),
    (Inches(3.05), "content.js",          ACCENT_BLUE,  "Scrapes DOM · Injects button & overlay"),
    (Inches(5.7),  "background.js\n(Service Worker)", ACCENT_PURP, "Orchestrates the pipeline · Holds ToolRegistry"),
    (Inches(8.35), "Tool Pipeline\n(5 Tools)", ACCENT_GREEN, "Runs ML + link scraping + Gemini LLM"),
    (Inches(11.0), "Result\nOverlay",     ACCENT_AMBER, "Verdict rendered on LinkedIn page"),
]
bw = Inches(2.3); bh = Inches(1.15)
by = Inches(1.8)
for bx, label, col, desc in boxes:
    add_rect(sl, bx, by, bw, bh, fill=col)
    add_text(sl, label, bx, by + Inches(0.12), bw, Inches(0.7),
             size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    # Arrow
    if bx < Inches(11.0):
        add_text(sl, "→", bx + bw, by + Inches(0.35), Inches(0.35), Inches(0.5),
                 size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(sl, desc, bx, by + bh + Inches(0.08), bw, Inches(0.45),
             size=9, color=MUTED, align=PP_ALIGN.CENTER)

# 5 Tools breakdown
add_text(sl, "The 5 Tools Inside background.js", Inches(0.45), Inches(3.5), Inches(12), Inches(0.4),
         size=15, bold=True, color=ACCENT_BLUE)

tools_v1 = [
    ("🤖  RoBERTaTool",           "ACCENT_GREEN", "POST to HuggingFace Inference API\naditya963/fraud-job-classifier\nReturns fraud probability 0→1"),
    ("🔗  DetectLinksTool",       "ACCENT_BLUE",  "Regex scan of job description\nFinds company URLs, social links\nQueues them for scraping"),
    ("🌐  LinkScraperTool",       "ACCENT_BLUE",  "Fetches each discovered URL\nExtracts visible text content\nBuilds evidence context"),
    ("📄  ContentAggregatorTool", "ACCENT_PURP",  "Merges all scraped text blobs\nCreates unified context string\nPassed to LLM as evidence"),
    ("🧠  JobAnalyzerTool",       "ACCENT_AMBER", "Calls Google Gemini 2.5-Flash\nReceives job text + RoBERTa + links\nOutputs verdict + reasoning"),
]
col_map = {"ACCENT_GREEN": ACCENT_GREEN, "ACCENT_BLUE": ACCENT_BLUE,
           "ACCENT_PURP": ACCENT_PURP, "ACCENT_AMBER": ACCENT_AMBER}

tw = Inches(2.45)
for i, (name, ckey, desc) in enumerate(tools_v1):
    tx = Inches(0.45) + i * Inches(2.57)
    ty = Inches(3.95)
    card(sl, tx, ty, tw, Inches(2.8))
    add_rect(sl, tx, ty, tw, Inches(0.06), fill=col_map[ckey])
    add_text(sl, name, tx + Inches(0.1), ty + Inches(0.12), tw - Inches(0.15), Inches(0.45),
             size=11, bold=True, color=col_map[ckey])
    add_text(sl, desc, tx + Inches(0.1), ty + Inches(0.58), tw - Inches(0.15), Inches(2.1),
             size=10, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — V1 DATA FLOW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "Version 1 — Data Flow",
            "Step-by-step: how a job listing becomes a fraud verdict")

tag(sl, "V1", Inches(12.3), Inches(0.2), w=Inches(0.75), h=Inches(0.3), color=ACCENT_GREEN)

steps = [
    (ACCENT_BLUE,  "① Scrape",
     "content.js reads the LinkedIn DOM\nExtracts: title, company, description,\nlocation, salary, links\nMutationObserver re-triggers on SPA navigation"),
    (ACCENT_PURP,  "② Message",
     "chrome.runtime.sendMessage()\n{ type: 'ANALYZE_JOB', data: jobData }\nKeeps channel open with  return true\nfor async response"),
    (ACCENT_GREEN, "③ Parallel Run",
     "Two branches simultaneously:\nBranch A → RoBERTaTool (HuggingFace)\nBranch B → DetectLinks → Scrape URLs\n         → Aggregate all content"),
    (ACCENT_AMBER, "④ LLM Verdict",
     "JobAnalyzerTool sends to Gemini:\n• Job text\n• RoBERTa probability\n• Scraped link evidence\nGemini returns: SAFE / SUSPICIOUS / LIKELY_FAKE"),
    (ACCENT_RED,   "⑤ Render",
     "background.js → chrome.tabs.sendMessage\ncontent.js injects overlay panel\nShows: verdict badge + probability\n+ LLM reasoning report"),
]

bw = Inches(2.38); bh = Inches(4.8); gap = Inches(0.2)
for i, (col, hdr, body) in enumerate(steps):
    bx = Inches(0.35) + i * (bw + gap)
    card(sl, bx, Inches(1.45), bw, bh)
    add_rect(sl, bx, Inches(1.45), bw, Inches(0.06), fill=col)
    add_text(sl, hdr, bx + Inches(0.12), Inches(1.56), bw - Inches(0.2), Inches(0.4),
             size=14, bold=True, color=col)
    add_text(sl, body, bx + Inches(0.12), Inches(2.05), bw - Inches(0.18), Inches(4.0),
             size=11, color=MUTED)
    if i < 4:
        add_text(sl, "▶", bx + bw + Inches(0.02), Inches(3.6), Inches(0.22), Inches(0.4),
                 size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Key note
add_rect(sl, Inches(0.35), Inches(6.5), Inches(12.6), Inches(0.65), fill=BG_CARD, line=DIVIDER, line_w=Pt(1))
add_text(sl, "⚡  Both RoBERTa and link scraping run in parallel (Promise.allSettled) — "
         "neither waits for the other.  Gemini only starts after both complete.",
         Inches(0.55), Inches(6.56), Inches(12.2), Inches(0.5),
         size=12, color=ACCENT_AMBER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MCP TOOLS WE BUILT (v1)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "MCP Tools We Built — v1 Tool Framework",
            "A custom tool system inspired by LangChain/MCP: every tool follows BaseTool → ToolResult → ToolRegistry")

tag(sl, "V1", Inches(12.3), Inches(0.2), w=Inches(0.75), h=Inches(0.3), color=ACCENT_GREEN)

# ── Top: three core abstractions ─────────────────────────────────────────────
abstractions = [
    (ACCENT_BLUE,  "BaseTool",
     "Abstract base class every tool extends\n"
     "Defines: name, description, run(input)\n"
     "Enforces a consistent interface so the\n"
     "pipeline can call any tool the same way"),
    (ACCENT_GREEN, "ToolResult",
     "Standardised return wrapper\n"
     "Fields: success, data, error, metadata\n"
     "ToolResult.ok(data)  — success path\n"
     "ToolResult.fail(msg) — error path\n"
     "Lets pipeline track success/failure\n"
     "and timing for every tool call"),
    (ACCENT_PURP,  "ToolRegistry",
     "Central registry for all tools\n"
     "registry.register(tool, category)\n"
     "registry.get(name) → tool instance\n"
     "registry.listTools() → all registered\n"
     "Modelled after LangChain's tool store\n"
     "but built for browser ES modules"),
    (ACCENT_AMBER, "PipelineBuilder",
     "Chains tools into a sequence\n"
     "Supports parallel branches\n"
     "Promise.allSettled for parallel runs\n"
     "Feeds output of one tool into next\n"
     "Controlled by PipelineConfig flags\n"
     "(e.g. skipLinkScraping = true)"),
]

aw = Inches(3.05)
for i, (col, hdr, body) in enumerate(abstractions):
    ax = Inches(0.35) + i * Inches(3.22)
    card(sl, ax, Inches(1.42), aw, Inches(3.5))
    add_rect(sl, ax, Inches(1.42), aw, Inches(0.06), fill=col)
    add_text(sl, hdr, ax + Inches(0.12), Inches(1.55), aw - Inches(0.18), Inches(0.42),
             size=14, bold=True, color=col)
    add_text(sl, body, ax + Inches(0.12), Inches(2.05), aw - Inches(0.18), Inches(2.75),
             size=10.5, color=MUTED)

# ── Bottom: the 5 actual tools built ─────────────────────────────────────────
add_text(sl, "The 5 Tools Built and Registered", Inches(0.45), Inches(5.15),
         Inches(12), Inches(0.4), size=14, bold=True, color=ACCENT_BLUE)

tools_built = [
    (ACCENT_GREEN, "RoBERTaTool",
     "Extends BaseTool\n"
     "Calls HuggingFace Inference API\n"
     "Model: aditya963/fraud-job-classifier\n"
     "Builds [SEP]-delimited input text\n"
     "Returns: fraud_probability + verdict"),
    (ACCENT_BLUE,  "DetectLinksTool",
     "Regex scan of job description\n"
     "Finds: company URLs, social links\n"
     "Filters: known CDNs + LinkedIn itself\n"
     "Returns: array of URLs to scrape"),
    (ACCENT_BLUE,  "LinkScraperTool",
     "fetch() each discovered URL\n"
     "Extracts visible text content\n"
     "Strips HTML, scripts, nav elements\n"
     "Returns: scraped text per URL"),
    (ACCENT_PURP,  "ContentAggregatorTool",
     "Merges all scraped text blobs\n"
     "Adds RoBERTa score to context\n"
     "Produces one unified string\n"
     "Input to the LLM analyzer"),
    (ACCENT_AMBER, "JobAnalyzerTool",
     "Calls Google Gemini 2.5-Flash\n"
     "System prompt: fraud detection expert\n"
     "Context: job text + score + links\n"
     "Returns: SAFE / SUSPICIOUS / FAKE\n"
     "+ full reasoning report"),
]

tw4 = Inches(2.45)
for i, (col, name, desc) in enumerate(tools_built):
    tx = Inches(0.35) + i * Inches(2.6)
    ty = Inches(5.6)
    card(sl, tx, ty, tw4, Inches(1.62))
    add_rect(sl, tx, ty, tw4, Inches(0.05), fill=col)
    add_text(sl, name, tx + Inches(0.1), ty + Inches(0.1), tw4 - Inches(0.15), Inches(0.36),
             size=11, bold=True, color=col)
    add_text(sl, desc, tx + Inches(0.1), ty + Inches(0.48), tw4 - Inches(0.15), Inches(1.1),
             size=9, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WHY VERSION 2
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "Why Build Version 2?",
            "v1 called Gemini directly with no tool calling — v2 introduced 13 purpose-built investigation tools via a deployed backend")

# Three columns: V1 Gap → What We Built → V2 Result
cols = [
    (ACCENT_RED,   "V1 — What Was Missing",
     ["Called Gemini directly with raw job text",
      "No structured tool calls — just one LLM prompt",
      "Gemini had to guess company legitimacy with no data",
      "No domain lookup, no WHOIS, no social check",
      "Only 5 internal tools, all browser-side",
      "No deployed model or backend infrastructure"]),
    (ACCENT_AMBER, "What We Built for V2",
     ["Deployed FastAPI backend on HuggingFace Spaces",
      "Built 13 specialised investigation tools server-side",
      "Deployed RoBERTa fraud classifier as a dedicated model",
      "Single batch API call runs all 13 tools in parallel",
      "LLM now receives structured evidence from every tool",
      "Backend shared with the web-app frontend"]),
    (ACCENT_GREEN, "V2 — What Improved",
     ["LLM verdict backed by real data: WHOIS, news, social",
      "13 tools vs 5 — far richer multi-signal coverage",
      "RoBERTa + domain + social + job boards + Wikipedia",
      "Weighted heuristic fallback if LLM unavailable",
      "Tool updates deploy via backend — no extension update",
      "Consistent analysis across extension and web-app"]),
]

cw = Inches(3.9)
for i, (col, hdr, items) in enumerate(cols):
    cx = Inches(0.4) + i * Inches(4.3)
    card(sl, cx, Inches(1.4), cw, Inches(5.7))
    add_rect(sl, cx, Inches(1.4), cw, Inches(0.07), fill=col)
    add_text(sl, hdr, cx + Inches(0.15), Inches(1.55), cw - Inches(0.2), Inches(0.45),
             size=14, bold=True, color=col)
    for j, item in enumerate(items):
        yy = Inches(2.15) + j * Inches(0.82)
        add_rect(sl, cx + Inches(0.15), yy + Inches(0.1), Inches(0.2), Inches(0.2), fill=col)
        add_text(sl, item, cx + Inches(0.5), yy + Inches(0.04), cw - Inches(0.58), Inches(0.65),
                 size=11, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — V2 ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "Version 2 — Architecture",
            "All API calls live in content.js · The service worker is a 5-line stub")

tag(sl, "V2", Inches(12.3), Inches(0.2), w=Inches(0.75), h=Inches(0.3), color=ACCENT_PURP)

# Deployed services at top right
card(sl, Inches(8.7), Inches(1.42), Inches(4.5), Inches(1.8))
add_text(sl, "Deployed Services", Inches(8.9), Inches(1.5), Inches(4.2), Inches(0.35),
         size=12, bold=True, color=ACCENT_BLUE)
services = [
    ("Backend API",  "hrmhrmhrm-company-backend-api.hf.space", ACCENT_BLUE),
    ("RoBERTa Model","hrmhrmhrm-roberta-model.hf.space",       ACCENT_GREEN),
    ("Frontend App", "hrmhrmhrm-company-frontend-app.hf.space",ACCENT_PURP),
]
for i, (name, url, col) in enumerate(services):
    yy = Inches(1.88) + i * Inches(0.42)
    add_text(sl, f"● {name}", Inches(8.9), yy, Inches(1.5), Inches(0.35), size=10, bold=True, color=col)
    add_text(sl, url, Inches(10.5), yy, Inches(2.5), Inches(0.35), size=9, color=MUTED)

# Main flow diagram
flow = [
    (Inches(0.3),  Inches(3.1), Inches(2.0), Inches(1.8), MUTED,        "LinkedIn\nJob Page",    "User opens job"),
    (Inches(2.9),  Inches(2.6), Inches(2.3), Inches(2.8), ACCENT_BLUE,  "content.js\n(All Logic)","Scrape · Call API · Render"),
    (Inches(5.85), Inches(1.55),Inches(2.4), Inches(1.1), MUTED,        "background.js",          "Ping stub only"),
    (Inches(5.85), Inches(3.1), Inches(2.4), Inches(1.0), ACCENT_PURP,  "/health",                "Wake HF Space"),
    (Inches(5.85), Inches(4.3), Inches(2.4), Inches(1.0), ACCENT_GREEN, "/api/v1/run-batch",      "13 tools, 90s timeout"),
    (Inches(5.85), Inches(5.5), Inches(2.4), Inches(1.0), ACCENT_AMBER, "/api/v1/llm/final-summary","LLM verdict, 50s timeout"),
    (Inches(9.0),  Inches(3.1), Inches(2.3), Inches(1.8), ACCENT_RED,   "FastAPI Backend",         "HuggingFace Spaces"),
    (Inches(11.6), Inches(3.1), Inches(1.5), Inches(1.8), ACCENT_GREEN, "Result\nOverlay",         "Rendered on page"),
]

for bx, by, bw2, bh2, col, label, hint in flow:
    add_rect(sl, bx, by, bw2, bh2, fill=col if col != MUTED else BG_CARD,
             line=col if col == MUTED else None, line_w=Pt(1))
    add_text(sl, label, bx + Inches(0.08), by + Inches(0.1), bw2 - Inches(0.12), Inches(0.75),
             size=12, bold=True, color=BG_DARK if col not in (MUTED, BG_CARD) else WHITE,
             align=PP_ALIGN.CENTER)
    add_text(sl, hint, bx + Inches(0.05), by + bh2 - Inches(0.35), bw2 - Inches(0.1), Inches(0.35),
             size=9, color=BG_DARK if col not in (MUTED, BG_CARD) else MUTED,
             align=PP_ALIGN.CENTER)

# Arrows
arrows = [
    (Inches(2.35), Inches(3.95), "→"),
    (Inches(8.45), Inches(3.95), "→"),
    (Inches(11.1), Inches(3.95), "→"),
]
for ax, ay, sym in arrows:
    add_text(sl, sym, ax, ay, Inches(0.4), Inches(0.45), size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

add_text(sl, "↑ not used for API calls",
         Inches(5.85), Inches(2.68), Inches(2.4), Inches(0.32), size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — V2 DATA FLOW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "Version 2 — Step-by-Step Data Flow",
            "From clicking 'FraudGuard' to seeing the verdict — 4 stages")

tag(sl, "V2", Inches(12.3), Inches(0.2), w=Inches(0.75), h=Inches(0.3), color=ACCENT_PURP)

stages = [
    (ACCENT_BLUE,  "Step 1\nScrape", Inches(0.35),
     "MutationObserver detects job panel\n"
     "Multi-selector fallback for each field:\n"
     "  • Title → 5 CSS selectors → all <h1>\n"
     "  • Company → selectors → /company/ links\n"
     "  • Description → containers → TreeWalker\n"
     "Retries up to 8× (500ms apart) for lazy\n"
     "loaded LinkedIn SPA content"),
    (ACCENT_PURP, "Step 2\nHealth Check", Inches(3.6),
     "GET /health  (12s timeout)\n"
     "If fails → retry with 50s timeout\n"
     "(HuggingFace Spaces sleep when idle)\n"
     "If both fail → user sees error:\n"
     "  'Cannot reach backend — verify URL'\n"
     "Confirms backend is awake before\n"
     "sending heavy batch request"),
    (ACCENT_GREEN, "Step 3\nBatch Tools", Inches(6.85),
     "POST /api/v1/run-batch  (90s timeout)\n"
     "Body: array of 13 tool requests\n"
     "Backend runs all 13 in parallel\n"
     "Returns: { ok, results[] }\n"
     "If batch fails → fallback to individual\n"
     "  scam_signals + roberta_classifier\n"
     "batchToDict() maps results by tool name"),
    (ACCENT_AMBER, "Step 4\nVerdict", Inches(10.1),
     "POST /api/v1/llm/final-summary (50s)\n"
     "LLM sees all 13 tool outputs + job data\n"
     "Returns: verdict + narrative report\n"
     "If LLM fails → heuristicVerdict():\n"
     "  Weighted score across all 13 tools\n"
     "  ≥55 pts = LIKELY_FAKE\n"
     "  ≥30 pts = SUSPICIOUS\n"
     "  <30 pts = SAFE"),
]

for col, hdr, bx, body in stages:
    bw3 = Inches(3.0); bh3 = Inches(5.55)
    card(sl, bx, Inches(1.42), bw3, bh3)
    add_rect(sl, bx, Inches(1.42), bw3, Inches(0.07), fill=col)
    add_text(sl, hdr, bx + Inches(0.1), Inches(1.52), bw3 - Inches(0.15), Inches(0.6),
             size=14, bold=True, color=col)
    add_text(sl, body, bx + Inches(0.1), Inches(2.2), bw3 - Inches(0.15), Inches(4.6),
             size=10.5, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THE 13 TOOLS (Part A: 1-7)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "The 13 Investigation Tools  (1 of 2)",
            "Sent as a single batch to the backend — all run in parallel server-side")

tag(sl, "V2", Inches(12.3), Inches(0.2), w=Inches(0.75), h=Inches(0.3), color=ACCENT_PURP)

tools_a = [
    (1, "roberta_classifier",    ACCENT_GREEN, "job_text",
     "Fine-tuned RoBERTa-base (125M params)\nDataset: EMSCAD (18k real/fake jobs)\nReturns fraud_probability 0→1\nThreshold: 0.87 = FRAUD\nWeight in heuristic: 40 pts max"),
    (2, "scam_signals",          ACCENT_RED,   "job_text",
     "Regex pattern library\nFlags: 'wire transfer', 'advance fee',\n'no experience needed', 'work from home'\nReturns: risk_level + scam_score/100\nWeight in heuristic: 25 pts max"),
    (3, "domain_reputation",     ACCENT_AMBER, "company website",
     "WHOIS lookup for company domain\nChecks: domain age, registrar, DNS flags\n<30 days old = HIGH risk\n<180 days old = MEDIUM risk\nWeight in heuristic: 15 pts max"),
    (4, "website_verify",        ACCENT_BLUE,  "company URL",
     "HTTP GET to company website\nChecks: site is live, SSL certificate valid\nNo SSL = suspicious signal\nSite down = strong fraud signal\nWeight in heuristic: 8 pts max"),
    (5, "social_profiles",       ACCENT_PURP,  "company_name",
     "Searches LinkedIn, Twitter, Facebook\nCounts number of platforms found\n0 platforms = strong fraud signal\n1 platform = weak fraud signal\nWeight in heuristic: 8 pts max"),
    (6, "job_boards",            ACCENT_GREEN, "title + company",
     "Cross-references Indeed, Glassdoor\nChecks if same job exists elsewhere\nJob not on any boards = suspicious\nCross-posting confirms legitimacy\nWeight in heuristic: 6 pts max"),
    (7, "company_wikipedia",     ACCENT_BLUE,  "company_name",
     "Wikipedia API search\nLooks for company article + extract\nNo Wikipedia = legitimacy concern\nFound + extract = positive signal\nWeight in heuristic: 4 pts max"),
]

tw2 = Inches(1.79)
for i, (num, name, col, inp, desc) in enumerate(tools_a):
    tx = Inches(0.3) + i * Inches(1.86)
    ty = Inches(1.45)
    card(sl, tx, ty, tw2, Inches(5.65))
    add_rect(sl, tx, ty, tw2, Inches(0.06), fill=col)
    add_text(sl, f"#{num}", tx + Inches(0.08), ty + Inches(0.1), Inches(0.35), Inches(0.32),
             size=10, bold=True, color=col)
    add_text(sl, name.replace("_", "_\n"), tx + Inches(0.08), ty + Inches(0.38),
             tw2 - Inches(0.15), Inches(0.7), size=10, bold=True, color=WHITE)
    add_text(sl, f"Input: {inp}", tx + Inches(0.08), ty + Inches(1.1),
             tw2 - Inches(0.15), Inches(0.32), size=9, italic=True, color=ACCENT_AMBER)
    add_text(sl, desc, tx + Inches(0.08), ty + Inches(1.45),
             tw2 - Inches(0.15), Inches(3.9), size=9.5, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — THE 13 TOOLS (Part B: 8-13)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "The 13 Investigation Tools  (2 of 2)",
            "Additional signals providing deeper context for the final verdict")

tag(sl, "V2", Inches(12.3), Inches(0.2), w=Inches(0.75), h=Inches(0.3), color=ACCENT_PURP)

tools_b = [
    (8,  "company_news",         ACCENT_AMBER, "company_name",
     "News API search\nCounts recent articles mentioning company\n0 articles = no media presence\nFound articles = legitimacy signal\nWeight in heuristic: 4 pts max"),
    (9,  "contact_info",         ACCENT_RED,   "job_text",
     "Scans for suspicious contact patterns\nFlags: personal Gmail, Yahoo, Hotmail\nFlags: WhatsApp-only contact requests\nFlags: foreign phone number formats\nSignal: recruiter credibility check"),
    (10, "salary_analysis",      ACCENT_GREEN, "job_text",
     "Extracts salary figures from text\nFlags unrealistic pay ranges\n'$5000/day no experience' = HIGH risk\nCompares against job title benchmarks\nSignal: bait-offer detection"),
    (11, "requirements_analysis",ACCENT_BLUE,  "job_text",
     "Analyzes job requirements section\nFlags vague or absent requirements\n'No experience needed' + high pay = fraud\nChecks for credential verification\nSignal: job posting quality"),
    (12, "company_registration",  ACCENT_PURP,  "company_name",
     "Business registry lookup\nChecks official company registrations\nVerifies company exists legally\nCross-checks incorporation date\nSignal: company legitimacy"),
    (13, "location_verify",      ACCENT_AMBER, "location + company",
     "Address and location plausibility\nChecks city/country combination\nVerifies company offices in claimed city\nFlags mismatches (e.g., 'Remote' abuse)\nSignal: geographic consistency"),
]

tw3 = Inches(2.05)
for i, (num, name, col, inp, desc) in enumerate(tools_b):
    tx = Inches(0.35) + i * Inches(2.16)
    ty = Inches(1.45)
    card(sl, tx, ty, tw3, Inches(5.65))
    add_rect(sl, tx, ty, tw3, Inches(0.06), fill=col)
    add_text(sl, f"#{num}", tx + Inches(0.08), ty + Inches(0.1), Inches(0.35), Inches(0.32),
             size=10, bold=True, color=col)
    add_text(sl, name.replace("_", "_\n"), tx + Inches(0.08), ty + Inches(0.38),
             tw3 - Inches(0.15), Inches(0.7), size=10, bold=True, color=WHITE)
    add_text(sl, f"Input: {inp}", tx + Inches(0.08), ty + Inches(1.1),
             tw3 - Inches(0.15), Inches(0.32), size=9, italic=True, color=ACCENT_AMBER)
    add_text(sl, desc, tx + Inches(0.08), ty + Inches(1.45),
             tw3 - Inches(0.15), Inches(3.9), size=9.5, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — VERDICT ENGINE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "Verdict Engine — How the Final Decision Is Made",
            "Two-path system: LLM narrative verdict + weighted heuristic fallback")

tag(sl, "V2", Inches(12.3), Inches(0.2), w=Inches(0.75), h=Inches(0.3), color=ACCENT_PURP)

# Left: LLM path
card(sl, Inches(0.4), Inches(1.42), Inches(5.9), Inches(5.65))
add_rect(sl, Inches(0.4), Inches(1.42), Inches(5.9), Inches(0.07), fill=ACCENT_AMBER)
add_text(sl, "Path A — LLM Final Summary  (primary)", Inches(0.6), Inches(1.55),
         Inches(5.5), Inches(0.45), size=13, bold=True, color=ACCENT_AMBER)
add_text(sl,
         "POST /api/v1/llm/final-summary\n\n"
         "Input:\n"
         "  • Full job posting (title, company, description)\n"
         "  • All 13 tool result objects\n"
         "  • Optional LLM config (key, model, base URL)\n\n"
         "The LLM reads every signal and writes:\n"
         "  • Human-readable fraud report\n"
         "  • Narrative explaining key red flags found\n"
         "  • Final verdict: SAFE / SUSPICIOUS / LIKELY_FAKE\n\n"
         "Default model: openai/gpt-4.1-mini via AIPipe\n"
         "Configurable: any OpenAI-compatible endpoint\n"
         "Timeout: 50 seconds",
         Inches(0.6), Inches(2.1), Inches(5.5), Inches(4.7),
         size=11, color=MUTED)

# Middle: OR / fallback arrow
add_text(sl, "If LLM\nfails ↓", Inches(6.5), Inches(3.5), Inches(0.8), Inches(0.8),
         size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Right: heuristic
card(sl, Inches(7.0), Inches(1.42), Inches(6.0), Inches(5.65))
add_rect(sl, Inches(7.0), Inches(1.42), Inches(6.0), Inches(0.07), fill=ACCENT_PURP)
add_text(sl, "Path B — Weighted Heuristic  (fallback)", Inches(7.2), Inches(1.55),
         Inches(5.6), Inches(0.45), size=13, bold=True, color=ACCENT_PURP)

weights = [
    ("roberta_classifier",    "40", ACCENT_GREEN),
    ("scam_signals",          "25", ACCENT_RED),
    ("domain_reputation",     "15", ACCENT_AMBER),
    ("website_verify",         "8", ACCENT_BLUE),
    ("social_profiles",        "8", ACCENT_BLUE),
    ("job_boards",             "6", MUTED),
    ("company_wikipedia",      "4", MUTED),
    ("company_news",           "4", MUTED),
]
add_text(sl, "Max pts  Tool", Inches(7.2), Inches(2.1), Inches(5.6), Inches(0.32),
         size=10, bold=True, color=MUTED)
add_rect(sl, Inches(7.2), Inches(2.42), Inches(5.65), Inches(0.012), fill=DIVIDER)
for i, (name, pts, col) in enumerate(weights):
    yy = Inches(2.5) + i * Inches(0.38)
    # pts bar
    bar_w = float(pts) / 40.0 * Inches(1.5)
    add_rect(sl, Inches(7.2), yy + Inches(0.05), bar_w, Inches(0.24), fill=col)
    add_text(sl, pts, Inches(7.2), yy, Inches(0.45), Inches(0.32),
             size=11, bold=True, color=col, align=PP_ALIGN.RIGHT)
    add_text(sl, name, Inches(7.75), yy, Inches(4.0), Inches(0.32),
             size=10, color=WHITE)

add_rect(sl, Inches(7.2), Inches(5.62), Inches(5.65), Inches(0.012), fill=DIVIDER)

verdicts_thresholds = [
    ("≥ 55 pts", "LIKELY_FAKE", ACCENT_RED),
    ("≥ 30 pts", "SUSPICIOUS",  ACCENT_AMBER),
    ("< 30 pts",  "SAFE",        ACCENT_GREEN),
]
for i, (thresh, label, col) in enumerate(verdicts_thresholds):
    tx = Inches(7.2) + i * Inches(2.0)
    add_rect(sl, tx, Inches(5.78), Inches(1.85), Inches(0.9), fill=col)
    add_text(sl, thresh, tx, Inches(5.82), Inches(1.85), Inches(0.35),
             size=10, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(sl, label, tx, Inches(6.15), Inches(1.85), Inches(0.35),
             size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — V1 vs V2 COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "Version 1 vs Version 2 — Side by Side",
            "A clear picture of what changed and why each decision was made")

rows = [
    ("Aspect",              "Version 1",                                  "Version 2",                                True),
    ("LLM approach",        "Direct Gemini call — raw job text only",     "LLM reads structured output of 13 tools",  False),
    ("Tool calling",        "No structured tool calls",                   "Batch API: all 13 tools run in parallel",  False),
    ("Backend",             "None — all logic inside the browser",        "FastAPI deployed on HuggingFace Spaces",   False),
    ("Number of tools",     "5 browser-side tools",                       "13 server-side investigation tools",       False),
    ("Company data",        "Gemini guesses from job text alone",         "WHOIS, social, news, job boards, Wikipedia",False),
    ("RoBERTa model",       "Called directly from browser (HF API)",      "Deployed as dedicated HF Space",           False),
    ("Verdict logic",       "RoBERTa threshold + one Gemini prompt",      "Weighted 13-signal heuristic + LLM report",False),
    ("LLM model",           "Google Gemini 2.5-Flash (fixed)",            "Any OpenAI-compatible model via AIPipe",   False),
    ("Tool updates",        "Requires re-installing extension",            "Deploy new backend — extension unchanged", False),
]

rh = Inches(0.42)
col_widths = [Inches(2.1), Inches(5.1), Inches(5.6)]
col_starts = [Inches(0.35), Inches(2.5), Inches(7.65)]

for i, row in enumerate(rows):
    is_header = row[3]
    aspect, v1, v2 = row[0], row[1], row[2]
    ry = Inches(1.42) + i * rh
    bg = BG_CARD if not is_header else ACCENT_BLUE
    add_rect(sl, Inches(0.35), ry, Inches(12.6), rh,
             fill=bg if is_header else (BG_CARD if i % 2 == 0 else BG_DARK),
             line=DIVIDER, line_w=Pt(0.5))
    texts = [aspect, v1, v2]
    for j, (cx, cw2, txt) in enumerate(zip(col_starts, col_widths, texts)):
        tcol = BG_DARK if is_header else (
            ACCENT_RED if j == 1 and not is_header and i > 0 and i in [1,4,5,6,8,9,10] else
            ACCENT_GREEN if j == 2 and not is_header and i > 0 else
            WHITE if is_header else MUTED if j == 0 else WHITE
        )
        add_text(sl, txt, cx + Inches(0.08), ry + Inches(0.04), cw2 - Inches(0.12), rh - Inches(0.06),
                 size=10.5 if not is_header else 12,
                 bold=is_header,
                 color=BG_DARK if is_header else tcol)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — KEY DESIGN DECISIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_title(sl, "Key Design Decisions",
            "The 'why' behind the technical choices — questions an examiner will ask")

decisions = [
    (ACCENT_BLUE,  "Why did v1 use direct Gemini and not tool calling?",
     "v1 was built before the backend was deployed. With no server infrastructure, the only option was "
     "to send the raw job text directly to Gemini and ask it to reason about fraud. "
     "This gave weak results because Gemini had no actual company data — just the text the scammer wrote."),
    (ACCENT_GREEN, "Why a custom 'LangChain-inspired' framework in v1?",
     "Real LangChain requires Node.js and cannot run inside a Chrome extension's browser context. "
     "The custom ToolRegistry + PipelineBuilder replicates the same tool-chaining pattern using "
     "pure ES modules that work in a service worker."),
    (ACCENT_AMBER, "Why deploy on HuggingFace Spaces?",
     "HuggingFace Spaces gives free GPU-backed hosting for ML models and free CPU hosting for APIs. "
     "The RoBERTa model (125M params) needs more compute than a browser can provide. "
     "Spaces also allows the backend to be shared between the web-app and the extension."),
    (ACCENT_PURP,  "Why a weighted heuristic instead of just RoBERTa?",
     "RoBERTa was trained on job text only — it cannot see domain age, social presence, or news coverage. "
     "A job posting can look textually legitimate while the company itself is fake. "
     "The weighted heuristic combines all 13 signals for a much harder-to-fool verdict."),
    (ACCENT_RED,   "Why build tools as a server-side batch API in v2?",
     "In v1, tools ran inside the browser — limited by the extension sandbox, no persistent state, "
     "and no access to server credentials. Moving tools to FastAPI means they can call paid APIs "
     "securely, run heavier ML models, and be updated without reinstalling the extension."),
    (ACCENT_GREEN, "Why move API keys to the server?",
     "Storing API keys in chrome.storage.local means any script running in the browser extension "
     "context could read them. Server-side environment variables are only accessible to the backend "
     "process — users only provide keys via the popup if they want to override the server default."),
]

dw = Inches(6.1); dh = Inches(1.6)
for i, (col, q, a) in enumerate(decisions):
    row = i // 2; ci = i % 2
    dx = Inches(0.35) + ci * Inches(6.49)
    dy = Inches(1.42) + row * Inches(1.85)
    card(sl, dx, dy, dw, dh)
    add_rect(sl, dx, dy, Inches(0.06), dh, fill=col)
    add_text(sl, q, dx + Inches(0.18), dy + Inches(0.1), dw - Inches(0.25), Inches(0.38),
             size=12, bold=True, color=col)
    add_text(sl, a, dx + Inches(0.18), dy + Inches(0.5), dw - Inches(0.25), Inches(1.0),
             size=10, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
add_rect(sl, 0, 0, W, Inches(0.08), fill=ACCENT_BLUE)
add_rect(sl, 0, Inches(0.08), W, Inches(0.04), fill=ACCENT_PURP)

add_text(sl, "Summary", Inches(1), Inches(0.6), Inches(11.3), Inches(0.7),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "FraudGuard Chrome Extension — Evolution from v1 to v2",
         Inches(1), Inches(1.3), Inches(11.3), Inches(0.45),
         size=16, color=MUTED, align=PP_ALIGN.CENTER)

points = [
    (ACCENT_BLUE,  "v1",
     "Self-contained browser extension · LangChain-inspired tool pipeline · "
     "5 custom tools · Direct Gemini API call · RoBERTa via HuggingFace"),
    (ACCENT_GREEN, "Why v2",
     "v1 sent raw job text to Gemini with no real company data · "
     "No tool calling meant Gemini was guessing · "
     "Needed structured signals: WHOIS, social, news, job boards"),
    (ACCENT_PURP,  "v2 Backend",
     "FastAPI deployed on HuggingFace Spaces · 13 investigation tools "
     "run server-side in parallel · RoBERTa model on dedicated HF Space"),
    (ACCENT_AMBER, "Verdict",
     "LLM now reads structured output from all 13 tools and writes a report · "
     "Weighted heuristic fallback if LLM unavailable · "
     "Score: ≥55=FAKE · ≥30=SUSPICIOUS · <30=SAFE"),
    (ACCENT_RED,   "Key Win",
     "Every verdict is backed by real data — not just what the scammer wrote · "
     "13 independent signals make the system far harder to fool"),
]

for i, (col, label, text) in enumerate(points):
    py = Inches(1.95) + i * Inches(0.98)
    add_rect(sl, Inches(0.6), py, Inches(0.85), Inches(0.6), fill=col)
    add_text(sl, label, Inches(0.6), py + Inches(0.1), Inches(0.85), Inches(0.4),
             size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(sl, text, Inches(1.65), py + Inches(0.08), Inches(11.0), Inches(0.72),
             size=12, color=MUTED)

add_rect(sl, 0, Inches(7.15), W, Inches(0.35), fill=BG_CARD)
add_text(sl, "Group 9  ·  DS & AI Lab Project  ·  FraudGuard",
         Inches(0.3), Inches(7.17), Inches(12.7), Inches(0.3),
         size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = "/Users/vishwas/Desktop/Webetention/Group-9-DS-and-AI-Lab-Project/webextension.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
